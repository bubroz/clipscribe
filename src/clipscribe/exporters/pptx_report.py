"""PowerPoint presentation generator for ClipScribe intelligence extraction.

Generates professional PPTX that opens in:
- Google Slides
- Microsoft PowerPoint
- Apple Keynote
"""

from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


def generate_pptx_report(
    transcript_result,
    intelligence_result,
    output_path: Path,
    filename: str = "executive_summary.pptx",
    is_series: bool = False,
    series_data: dict = None,
):
    """Generate professional PowerPoint presentation.

    Args:
        transcript_result: TranscriptResult from provider
        intelligence_result: IntelligenceResult from provider
        output_path: Directory to save presentation
        filename: Output filename
        is_series: Whether this is a series analysis
        series_data: Aggregate data for series (if is_series=True)

    Returns:
        Path to generated PPTX file
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    if is_series and series_data:
        title.text = "Series Intelligence Analysis"
        subtitle.text = f"{series_data.get('series_name', 'Analysis')}\n{series_data.get('total_videos', 0)} Videos Processed\n{datetime.now().strftime('%B %d, %Y')}"
    else:
        title.text = "Intelligence Extraction Report"
        subtitle.text = f"ClipScribe Analysis\n{datetime.now().strftime('%B %d, %Y')}"

    # Slide 2: Executive Summary
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Executive Summary"

    body = slide.placeholders[1].text_frame
    body.clear()

    if is_series and series_data:
        # Series metrics
        p = body.paragraphs[0]
        p.text = f"Processed {series_data.get('total_videos', 0)} videos"
        p.level = 0

        p = body.add_paragraph()
        p.text = f"{series_data.get('unique_entities', 0)} unique entities across corpus"
        p.level = 1

        p = body.add_paragraph()
        p.text = f"Top entity appeared in {series_data.get('max_entity_frequency', 0)} videos"
        p.level = 1
    else:
        # Single video metrics
        p = body.paragraphs[0]
        p.text = f"{len(intelligence_result.entities)} entities extracted"
        p.level = 0

        p = body.add_paragraph()
        p.text = f"{len(intelligence_result.relationships)} relationships mapped"
        p.level = 1

        p = body.add_paragraph()
        p.text = f"{transcript_result.speakers} speakers identified"
        p.level = 1

        p = body.add_paragraph()
        p.text = f"{len(intelligence_result.topics)} main topics"
        p.level = 1

        p = body.add_paragraph()
        p.text = f"Processing cost: ${transcript_result.cost + intelligence_result.cost:.4f}"
        p.level = 1

    # Slide 3: Top Entities
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Key Entities"

    body = slide.placeholders[1].text_frame
    body.clear()

    entities_to_show = intelligence_result.entities[:10] if intelligence_result.entities else []
    for entity in entities_to_show:
        p = body.add_paragraph() if body.paragraphs else body.paragraphs[0]
        name = entity.get("name", "Unknown")
        type_val = entity.get("type", "UNKNOWN")
        conf = entity.get("confidence", 0)
        p.text = f"{name} ({type_val}) - Confidence: {conf:.2f}"
        p.level = 0

    # Slide 4: Relationships
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Key Relationships"

    body = slide.placeholders[1].text_frame
    body.clear()

    rels_to_show = (
        intelligence_result.relationships[:8] if intelligence_result.relationships else []
    )
    for rel in rels_to_show:
        p = body.add_paragraph() if body.paragraphs else body.paragraphs[0]
        subject = rel.get("subject", "Unknown")
        predicate = rel.get("predicate", "related_to")
        obj = rel.get("object", "Unknown")
        p.text = f"{subject} → {predicate} → {obj}"
        p.level = 0

    # Slide 5: Topics & Timeline
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Topics & Timeline"

    body = slide.placeholders[1].text_frame
    body.clear()

    for topic in intelligence_result.topics:
        p = body.add_paragraph() if body.paragraphs else body.paragraphs[0]
        name = topic.get("name", "Unknown")
        relevance = topic.get("relevance", 0)
        time_range = topic.get("time_range", "N/A")
        p.text = f"{name} ({time_range}) - Relevance: {relevance:.2f}"
        p.level = 0

    # Slide 6: Key Moments
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Key Moments"

    body = slide.placeholders[1].text_frame
    body.clear()

    for moment in intelligence_result.key_moments[:6]:
        p = body.add_paragraph() if body.paragraphs else body.paragraphs[0]
        timestamp = moment.get("timestamp", "00:00")
        desc = moment.get("description", "No description")
        sig = moment.get("significance", 0)
        p.text = f"{timestamp} - {desc} (significance: {sig:.2f})"
        p.level = 0

    # Slide 7: Sentiment
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Sentiment Analysis"

    body = slide.placeholders[1].text_frame
    body.clear()

    sentiment = intelligence_result.sentiment
    p = body.paragraphs[0]
    p.text = f"Overall: {sentiment.get('overall', 'unknown').upper()}"
    p.level = 0

    if sentiment.get("per_topic"):
        p = body.add_paragraph()
        p.text = "Per-Topic Breakdown:"
        p.level = 1

        for topic, topic_sent in list(sentiment.get("per_topic", {}).items())[:5]:
            p = body.add_paragraph()
            p.text = f"{topic}: {topic_sent}"
            p.level = 2

    # Save
    full_path = output_path / filename
    prs.save(str(full_path))

    return full_path
